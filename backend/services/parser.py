from typing import List, Dict, Any
import fitz # pymupdf
from pptx import Presentation
import os

def extract_text_from_pdf(file_path: str) -> List[Dict[str, Any]]:
    """
    Extracts text from a PDF file.
    Returns a list of dicts: {'page': int, 'text': str}
    """
    results = []
    try:
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                results.append({"page": i + 1, "text": text.strip()})
        return results
    except Exception as e:
        print(f"Error extracting PDF {file_path}: {e}")
        return []

def extract_text_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extracts text from a PPTX file.
    Returns a list of dicts: {'page': int, 'text': str}
    """
    results = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)

            full_text = "\n".join(text_runs)
            if full_text.strip():
                results.append({"page": i + 1, "text": full_text.strip()})
        return results
    except Exception as e:
        print(f"Error extracting PPTX {file_path}: {e}")
        return []

def extract_text(file_path: str, content_type: str) -> List[Dict[str, Any]]:
    if content_type == "application/pdf":
        return extract_text_from_pdf(file_path)
    elif content_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
        return extract_text_from_pptx(file_path)
    else:
        # Fallback or unknown
        return []
